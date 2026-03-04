from __future__ import annotations

from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation

from utils import read_json_file, read_text_file


class TextExtractor:
    def __init__(self, max_file_size_mb: int = 25) -> None:
        self.max_file_size_bytes = max_file_size_mb * 1024 * 1024

    def extract_text(self, path: Path) -> str:
        if path.stat().st_size > self.max_file_size_bytes:
            raise ValueError(f"File exceeds size limit ({self.max_file_size_bytes} bytes): {path}")

        extension = path.suffix.lower()
        if extension == ".pdf":
            return self._extract_pdf(path)
        if extension == ".docx":
            return self._extract_docx(path)
        if extension == ".pptx":
            return self._extract_pptx(path)
        if extension == ".json":
            return read_json_file(path)
        return read_text_file(path)

    def _extract_pdf(self, path: Path) -> str:
        parts: list[str] = []
        with fitz.open(path) as document:
            for page in document:
                parts.append(page.get_text("text"))
        return "\n".join(parts)

    def _extract_docx(self, path: Path) -> str:
        document = Document(str(path))
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
        return "\n".join(parts)

    def _extract_pptx(self, path: Path) -> str:
        presentation = Presentation(str(path))
        parts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if text:
                    parts.append(text)
        return "\n".join(parts)
